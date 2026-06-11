"""
Génère la présentation finale (PowerPoint) du projet DAC.
Sections imposées : problème traité, solution proposée, choix techniques,
démonstration fonctionnelle, limites restantes, améliorations possibles.

Usage : python scripts/generate_pptx.py
Sortie : presentation-dac.pptx (racine du repo)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Palette
PRIMARY = RGBColor(0x63, 0x66, 0xF1)   # indigo
DARK = RGBColor(0x0F, 0x17, 0x2A)      # bleu nuit
GREEN = RGBColor(0x10, 0xB9, 0x81)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
GREY = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW, SH = prs.slide_width, prs.slide_height


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_box(slide, left, top, width, height):
    from pptx.enum.shapes import MSO_SHAPE
    return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)


def title_slide(title, subtitle, footer):
    slide = prs.slides.add_slide(BLANK)
    bg = add_box(slide, 0, 0, SW, SH)
    _fill(bg, DARK)
    # bande
    band = add_box(slide, 0, Inches(2.6), SW, Inches(1.7))
    _fill(band, PRIMARY)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE

    sb = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.0))
    sp = sb.text_frame.paragraphs[0]
    sr = sp.add_run(); sr.text = subtitle
    sr.font.size = Pt(22); sr.font.color.rgb = RGBColor(0xC7, 0xD2, 0xFE)

    fb = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
    fp = fb.text_frame.paragraphs[0]
    fr = fp.add_run(); fr.text = footer
    fr.font.size = Pt(14); fr.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    return slide


def content_slide(kicker, title, bullets):
    """bullets: list de (texte, niveau) ; niveau 0 = puce principale, 1 = sous-puce."""
    slide = prs.slides.add_slide(BLANK)
    bg = add_box(slide, 0, 0, SW, SH)
    _fill(bg, WHITE)
    # barre latérale
    side = add_box(slide, 0, 0, Inches(0.25), SH)
    _fill(side, PRIMARY)
    # kicker
    kb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.5))
    kp = kb.text_frame.paragraphs[0]
    kr = kp.add_run(); kr.text = kicker.upper()
    kr.font.size = Pt(14); kr.font.bold = True; kr.font.color.rgb = PRIMARY
    # titre
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.95), Inches(12), Inches(1.0))
    tp = tb.text_frame.paragraphs[0]
    tr = tp.add_run(); tr.text = title
    tr.font.size = Pt(32); tr.font.bold = True; tr.font.color.rgb = DARK
    # ligne
    ln = add_box(slide, Inches(0.7), Inches(1.95), Inches(11.9), Pt(2))
    _fill(ln, LIGHT)
    # contenu
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(4.9))
    tf = cb.text_frame
    tf.word_wrap = True
    first = True
    for text, lvl in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        run = p.add_run()
        run.text = ("• " if lvl == 0 else "– ") + text
        run.font.size = Pt(20 if lvl == 0 else 17)
        run.font.bold = lvl == 0
        run.font.color.rgb = DARK if lvl == 0 else GREY
        p.space_after = Pt(8)
    return slide


def section_break(num, title):
    slide = prs.slides.add_slide(BLANK)
    bg = add_box(slide, 0, 0, SW, SH)
    _fill(bg, PRIMARY)
    nb = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(12), Inches(1.2))
    np = nb.text_frame.paragraphs[0]
    nr = np.add_run(); nr.text = f"0{num}"
    nr.font.size = Pt(60); nr.font.bold = True; nr.font.color.rgb = RGBColor(0xC7, 0xD2, 0xFE)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(12), Inches(1.2))
    tp = tb.text_frame.paragraphs[0]
    tr = tp.add_run(); tr.text = title
    tr.font.size = Pt(40); tr.font.bold = True; tr.font.color.rgb = WHITE
    return slide


# ----------------------------------------------------------------------------
# 0. Titre
# ----------------------------------------------------------------------------
title_slide(
    "DAC — DevOps-as-a-Chat",
    "Améliorer l'expérience du chat & sécuriser les actions (dry-run + confirmation)",
    "CodeCamp ETNA 2026  ·  Challenges 1 & 2  ·  Présentation finale",
)

# ----------------------------------------------------------------------------
# 1. Problème traité
# ----------------------------------------------------------------------------
content_slide(
    "Section 1", "Le problème traité",
    [
        ("DAC permet de créer/configurer de l'infrastructure AWS via un chat.", 0),
        ("Challenge 1 — l'expérience du chat était confuse :", 0),
        ("Messages d'erreur bruts, tronqués, avec doublons (« 500: … 500: … »).", 1),
        ("Aucune distinction entre information, action proposée et action exécutée.", 1),
        ("Pas de feedback de progression : on ne sait pas si un déploiement se termine.", 1),
        ("Challenge 2 — les actions sensibles n'étaient pas sécurisées :", 0),
        ("La configuration s'exécutait dès la sélection des VM, sans confirmation.", 1),
        ("Aucune simulation possible avant exécution, aucune trace des décisions.", 1),
    ],
)

# ----------------------------------------------------------------------------
# 2. Solution proposée
# ----------------------------------------------------------------------------
content_slide(
    "Section 2", "La solution proposée",
    [
        ("Challenge 1 — un chat plus clair et plus lisible :", 0),
        ("Traduction des erreurs techniques → message clair + action corrective.", 1),
        ("Badges par type de message : info / proposition / exécution / erreur.", 1),
        ("Récapitulatif structuré du plan + barre de progression des tâches.", 1),
        ("Suggestions de saisie + messages d'aide compréhensibles.", 1),
        ("Challenge 2 — sécurisation des actions :", 0),
        ("Classification de sensibilité (safe / sensitive / dangerous).", 1),
        ("Mode simulation (dry-run) : « simuler <commande> » sans rien exécuter.", 1),
        ("Confirmation explicite + journalisation des décisions (oui/non).", 1),
    ],
)

# ----------------------------------------------------------------------------
# 3. Choix techniques
# ----------------------------------------------------------------------------
content_slide(
    "Section 3", "Les choix techniques",
    [
        ("Stack : FastAPI (Python) + React/TypeScript (Vite) + PostgreSQL, Docker Compose.", 0),
        ("Backend — services découplés et testables :", 0),
        ("error_translator, action_safety, plan_presenter, decision_log (table dédiée).", 1),
        ("AMI résolue via data source aws_ami ; dry-run branché dans l'exécuteur SSM.", 1),
        ("Frontend — composants MUI + rendu Markdown des messages :", 0),
        ("Type de message porté par extra.type, avec inférence de secours.", 1),
        ("Historique masqué au login, révélé au scroll (timestamp de connexion).", 1),
        ("Qualité : 17 tests frontend (Vitest + Testing Library) sur les briques clés.", 0),
    ],
)

# ----------------------------------------------------------------------------
# 4. Démonstration fonctionnelle
# ----------------------------------------------------------------------------
content_slide(
    "Section 4", "La démonstration fonctionnelle",
    [
        ("Scénario 1 — création d'infrastructure :", 0),
        ("« crée une instance ubuntu sur aws » → plan → confirmation → déploiement suivi.", 1),
        ("Scénario 2 — erreur claire :", 0),
        ("Type d'instance non Free Tier → message « utilise t3.micro » au lieu d'une stacktrace.", 1),
        ("Scénario 3 — simulation (dry-run) :", 0),
        ("« simuler sudo systemctl restart nginx » → commande affichée, AUCUNE exécution.", 1),
        ("Scénario 4 — confirmation d'une configuration :", 0),
        ("Sélection VM → plan d'action + boutons Confirmer/Annuler → décision tracée.", 1),
    ],
)

# ----------------------------------------------------------------------------
# 5. Limites restantes
# ----------------------------------------------------------------------------
content_slide(
    "Section 5", "Les limites restantes",
    [
        ("La commande shell exacte du flux configure n'est pas toujours affichée :", 0),
        ("elle est générée à l'exécution ; on montre l'action de catalogue.", 1),
        ("Journalisation des décisions branchée surtout sur create et configure.", 0),
        ("Statut des conversations (déployé / en erreur) dans la sidebar : non fait (déféré).", 0),
        ("Dry-run SSM : démontré via la commande « simuler », pas encore un toggle UI.", 0),
        ("Tests : couverture sur les briques pures ; pas encore de tests E2E.", 0),
    ],
)

# ----------------------------------------------------------------------------
# 6. Améliorations possibles
# ----------------------------------------------------------------------------
content_slide(
    "Section 6", "Les améliorations possibles",
    [
        ("Toggle « Simuler / Exécuter » dans l'UI + dry-run Ansible natif (--check).", 0),
        ("Écran d'historique des décisions (audit trail consultable).", 0),
        ("Badges de statut + reprise de conversation dans la sidebar (Challenge 1, Axe 4).", 0),
        ("Garde-fou systématique : refus d'exécution directe pour toute action sensible.", 0),
        ("Étendre la journalisation à audit/monitoring et la traduction d'erreurs partout.", 0),
        ("Tests E2E (Playwright) sur les parcours création / configuration / simulation.", 0),
    ],
)

# ----------------------------------------------------------------------------
# Clôture
# ----------------------------------------------------------------------------
title_slide(
    "Merci !",
    "DAC — un chat DevOps plus clair, et des actions sécurisées par simulation & confirmation.",
    "Questions ?  ·  CodeCamp ETNA 2026",
)

OUT = "presentation-dac.pptx"
prs.save(OUT)
print(f"OK -> {OUT} ({len(prs.slides._sldIdLst)} slides)")
